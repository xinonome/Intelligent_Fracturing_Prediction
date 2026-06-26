
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path

SRC = Path(r'C:\Workspace\Graph\project_report_4_23_source.pptx')
OUT = Path(r'C:\Workspace\Graph\project_report_4_23_gnn_added_v2.pptx')
FIG = Path(r'C:\Workspace\Graph\report_figures')
prs = Presentation(str(SRC))
TITLE = RGBColor(22, 92, 72); DARK = RGBColor(45,55,72); MUTED = RGBColor(95,105,115)
GREEN = RGBColor(46,125,94); ORANGE = RGBColor(242,140,40); RED = RGBColor(210,74,67); BLUE = RGBColor(55,116,184)
blank = prs.slide_layouts[0]

def set_tf(tf, text, size=16, color=DARK, bold=False, align=None):
    tf.clear(); p=tf.paragraphs[0]; p.text=text
    p.font.name='Microsoft YaHei'; p.font.size=Pt(size); p.font.color.rgb=color; p.font.bold=bold
    if align: p.alignment=align
    return p

def title(slide, text, sub=None):
    box=slide.shapes.add_textbox(Inches(.55), Inches(.28), Inches(12.1), Inches(.55)); set_tf(box.text_frame,text,24,TITLE,True)
    line=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.55), Inches(.88), Inches(11.9), Inches(.03)); line.fill.solid(); line.fill.fore_color.rgb=TITLE; line.line.color.rgb=TITLE
    if sub:
        b=slide.shapes.add_textbox(Inches(.58), Inches(.94), Inches(11.7), Inches(.35)); set_tf(b.text_frame,sub,12,MUTED)
    tag=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.65), Inches(.25), Inches(2.0), Inches(.38)); tag.fill.solid(); tag.fill.fore_color.rgb=RGBColor(232,244,238); tag.line.color.rgb=RGBColor(185,215,198)
    set_tf(tag.text_frame,'三、基于GNN的标签识别',10,TITLE,True,PP_ALIGN.CENTER)

def bullets(slide,x,y,w,h,items,size=14):
    box=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=box.text_frame; tf.clear()
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=it; p.font.name='Microsoft YaHei'; p.font.size=Pt(size); p.font.color.rgb=DARK; p.space_after=Pt(5)

def card(slide,x,y,w,h,value,label,color):
    shp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h)); shp.fill.solid(); shp.fill.fore_color.rgb=RGBColor(248,251,249); shp.line.color.rgb=RGBColor(200,220,210)
    tf=shp.text_frame; tf.clear(); p=tf.paragraphs[0]; p.text=value; p.font.name='Microsoft YaHei'; p.font.size=Pt(23); p.font.bold=True; p.font.color.rgb=color; p.alignment=PP_ALIGN.CENTER
    p2=tf.add_paragraph(); p2.text=label; p2.font.name='Microsoft YaHei'; p2.font.size=Pt(10); p2.font.color.rgb=MUTED; p2.alignment=PP_ALIGN.CENTER

def pic(slide,name,x,y,w,h):
    p=FIG/name
    if p.exists(): slide.shapes.add_picture(str(p),Inches(x),Inches(y),width=Inches(w),height=Inches(h))

def move_slide(prs, old, new):
    sldIdLst=prs.slides._sldIdLst; s=list(sldIdLst); el=s[old]; sldIdLst.remove(el); sldIdLst.insert(new,el)

def add_task_slide():
    s=prs.slides.add_slide(blank); title(s,'三、基于GNN的标签识别：任务定义与建模思路','以压裂段内连续秒点数据为输入，识别当前工况并预测下一工况转移概率')
    card(s,.75,1.35,2.2,1.0,'10 秒/点','采样粒度',BLUE); card(s,3.15,1.35,2.2,1.0,'1-6 点','输入窗口实验',GREEN); card(s,5.55,1.35,2.2,1.0,'WORKING_TYPE','预测标签',ORANGE); card(s,7.95,1.35,2.2,1.0,'GNN / LGBM','模型对照',RED)
    bullets(s,.8,2.65,5.7,3.2,['样本单位：每个压裂段是一条连续时序，窗口内相邻采样点构成链式图。','节点特征：砂比、泵压、排量、累计液量等施工参数及动态统计特征。','GNN 思路：窗口内每个时间点作为节点，通过时间边聚合局部上下文。','评估策略：按段划分训练/测试/验证/迁移集，避免相邻点泄漏。'],15)
    for i,(txt,col) in enumerate([('连续秒点',BLUE),('滑动窗口',GREEN),('链式图/GNN',ORANGE),('工况标签',RED)]):
        x=6.9+i*1.45; shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(3.05),Inches(1.25),Inches(.75)); shp.fill.solid(); shp.fill.fore_color.rgb=RGBColor(250,250,250); shp.line.color.rgb=col; set_tf(shp.text_frame,txt,12,col,True,PP_ALIGN.CENTER)
        if i<3:
            arr=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(x+1.18),Inches(3.27),Inches(.35),Inches(.28)); arr.fill.solid(); arr.fill.fore_color.rgb=RGBColor(180,190,190); arr.line.color.rgb=RGBColor(180,190,190)
    bullets(s,6.85,4.25,5.65,1.4,['新增严格转移模型：输入当前窗口 + 当前工况，输出下一点各工况概率。','适合做“当前状态下的下一步风险评估”。'],14)

def add_data_slide():
    s=prs.slides.add_slide(blank); title(s,'数据基础：新增数据显著改善异常样本覆盖','旧数据异常仅约 1.25%，新增数据异常约 23.16%，合并后异常提升到约 5.00%')
    pic(s,'01_normal_abnormal_ratio.png',.45,1.25,6.1,3.7); pic(s,'02_abnormal_label_distribution.png',6.65,1.25,6.0,3.7)
    bullets(s,.7,5.25,11.8,.9,['新增工况包括：缝内暂堵、缝高延伸、延伸受阻、滤失过大；为多分类识别和迁移学习提供了更丰富的异常样本。','仍存在长尾问题：滤失过大、延伸受阻、其他等类别样本量仍偏少，需要继续补充。'],13)

def add_model_slide():
    s=prs.slides.add_slide(blank); title(s,'模型路线迭代：从二分类识别到合并数据多分类','当前最优多分类测试集 Macro F1 = 0.5739，新增数据后多分类已具备可行性')
    pic(s,'04_model_macro_f1_comparison.png',.55,1.25,6.2,3.9); pic(s,'05_multiclass_per_class_metrics.png',6.85,1.25,5.8,3.9)
    bullets(s,.75,5.35,11.6,.75,['关键结论：性能提升主要来自“加砂后截断 + 动态特征 + 新增异常样本”，而不是单纯增加模型复杂度。','当前难点：部分类别在测试集覆盖不足，极少数类仍需更多跨段样本支撑泛化。'],13)

def add_transition_slide():
    s=prs.slides.add_slide(blank); title(s,'下一工况转移概率模型：从识别走向风险预警','严格预测 y(t+1)：输入当前窗口特征 + 当前工况，输出下一点各工况概率分布')
    card(s,.75,1.2,2.4,.9,'0.9846','下一工况预测 Test Macro F1',GREEN); card(s,3.35,1.2,2.4,.9,'99.9895%','正常→正常平均概率',BLUE); card(s,5.95,1.2,2.4,.9,'11.63%','单点最高异常风险示例',ORANGE)
    pic(s,'06_normal_to_abnormal_transition_probs.png',.65,2.45,5.8,3.2)
    bullets(s,6.85,2.55,5.55,3.1,['输出形式：P(下一点=正常/砂堵/缝内暂堵/缝口暂堵/...)。','解释边界：这是状态转移预测模型，依赖当前工况标签，不能与纯识别模型直接比较。','业务价值：可以在已知当前状态下，给出下一采样点发生异常转移的风险概率。','示例：当前正常时，部分点下一步转为缝内暂堵的预测概率最高可达 11.63%。'],14)

def add_transfer_slide():
    s=prs.slides.add_slide(blank); title(s,'迁移学习与后续支撑需求','少量目标域样本微调有效，但进一步提升依赖异常样本覆盖和标签质量')
    card(s,.8,1.2,2.55,1.0,'0.3122 → 0.3310','迁移集 Macro F1',GREEN); card(s,3.65,1.2,2.55,1.0,'+2.73%','迁移集 Accuracy 提升',BLUE); card(s,6.5,1.2,2.55,1.0,'0.5739 → 0.4389','删除纯正常段后下降',RED)
    bullets(s,.85,2.65,5.7,2.95,['迁移学习实验：源域预训练 GNN，迁移集 30% support 微调，70% query 评估。','结果：目标域 Macro F1 从 0.3122 提升到 0.3310，方向有效但提升有限。','纯正常段不能直接删除：它们仍对正常边界学习有价值，删除后效果明显下降。'],14)
    bullets(s,6.9,2.65,5.4,2.95,['需要甲方支持：补充更多含异常段，尤其是低频异常和跨井跨段样本。','统一标签口径：明确各工况定义、起止时间和空白标签处理规则。','提供业务上下文：施工阶段、人工处置、设计参数、专家复核结论。'],14)

original_count=len(prs.slides)
for f in [add_task_slide, add_data_slide, add_model_slide, add_transition_slide, add_transfer_slide]:
    f()
insert_at=13
for j in range(5):
    move_slide(prs, original_count + j, insert_at + j)
prs.save(str(OUT))
print(OUT)
print(len(prs.slides))
